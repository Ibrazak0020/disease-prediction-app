# ------------------------------------------------------------
# AI Disease Prediction Presentation Generator
# Author: Kazeem Ibrahim Opeyemi
# Institution: Abiola Ajimobi Technical University
# Supervisor: Dr. JET
# Session: 2024/2025
# ------------------------------------------------------------

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Theme colors
BLUE = RGBColor(0, 51, 102)
LIGHT_BLUE = RGBColor(0, 102, 204)

# Helper function
def add_slide(title, content, notes=""):
    slide_layout = prs.slide_layouts[5]  # blank layout
    slide = prs.slides.add_slide(slide_layout)

    # background gradient (simple effect)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(230, 240, 255)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = BLUE

    # Content
    bodyBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.5), Inches(4.5))
    tf2 = bodyBox.text_frame
    tf2.word_wrap = True
    for line in content:
        p = tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(20, 20, 20)

    # Footer
    footerBox = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(12), Inches(0.5))
    ft = footerBox.text_frame
    ft.text = "Abiola Ajimobi Technical University"
    ft.paragraphs[0].font.size = Pt(12)
    ft.paragraphs[0].font.color.rgb = BLUE
    ft.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add notes
    slide.notes_slide.notes_text_frame.text = notes


# ------------------------------------------------------------
# SLIDES
# ------------------------------------------------------------
add_slide(
    "1. AI Disease Prediction System Using Machine Learning and Streamlit",
    [
        "Student: Kazeem Ibrahim Opeyemi",
        "Supervisor: Dr. JET",
        "Institution: Abiola Ajimobi Technical University",
        "Session: 2024/2025",
        "",
        "Insert University Logo Here"
    ],
    "Opening slide introducing the project, supervisor, and session."
)

add_slide(
    "2. Introduction",
    [
        "This project leverages machine learning and Streamlit to predict diseases from symptoms.",
        "It uses historical symptom–disease datasets to train a classification model.",
        "The app provides users with likely diagnoses, descriptions, and precautions."
    ],
    "Briefly explain the aim: to assist early disease detection using AI."
)

add_slide(
    "3. Problem Statement",
    [
        "Healthcare facilities face challenges in providing quick, accessible diagnosis.",
        "Manual evaluation of symptoms can lead to human error.",
        "An automated system can provide faster preliminary assessments."
    ],
    "Mention that this system is not a replacement for doctors but a supportive tool."
)

add_slide(
    "4. Project Objectives",
    [
        "Develop a machine learning model to predict possible diseases.",
        "Design a Streamlit web application for interactive prediction.",
        "Integrate datasets for disease symptoms, descriptions, and precautions.",
        "Ensure accurate and user-friendly disease diagnosis."
    ],
    "Clearly state objectives as expected by supervisor."
)

add_slide(
    "5. Scope of the Project",
    [
        "Focuses on symptom-based disease prediction.",
        "Model trained on pre-defined datasets (disease.csv, symptom_description.csv, etc.).",
        "Output limited to top three likely diseases with confidence scores.",
        "Educational and research purpose — not medical advice."
    ],
    "Define project boundaries and purpose."
)

add_slide(
    "6. Methodology",
    [
        "1. Data Collection and Preprocessing.",
        "2. Feature Selection: Symptoms as binary features (0/1).",
        "3. Model Training: RandomForestClassifier for prediction.",
        "4. Evaluation: Achieved 1.0 accuracy on test data.",
        "5. Deployment: Streamlit web application."
    ],
    "Explain Random Forest model and Streamlit deployment."
)

add_slide(
    "7. System Design and Implementation",
    [
        "Architecture includes: Dataset → Model Training → Prediction → Web Interface.",
        "Frontend: Streamlit for user input and visualization.",
        "Backend: Scikit-learn model served via pickle.",
        "Outputs: Predicted disease, description, and precautions."
    ],
    "Discuss the workflow diagram and implementation details."
)

add_slide(
    "8. Results and Model Accuracy",
    [
        "Model achieved 100% accuracy (1.0) on test dataset.",
        "Confusion matrix showed no misclassifications.",
        "High accuracy attributed to balanced and well-structured data.",
        "Insert Model Accuracy Graph Placeholder Here."
    ],
    "Highlight your strong accuracy result confidently."
)

add_slide(
    "9. App Interface and Functionality",
    [
        "Users select symptoms via checkboxes.",
        "System predicts top three possible diseases with confidence scores.",
        "Displays disease description and recommended precautions.",
        "Insert App Screenshot Placeholder Here."
    ],
    "Show your app demo here."
)

add_slide(
    "10. Conclusion and Recommendations",
    [
        "The project demonstrates that machine learning can enhance preliminary disease diagnosis.",
        "Encouraged to expand dataset for broader coverage.",
        "Future improvement: include deep learning and multilingual support.",
        "",
        "Thank you for your time and attention.",
        "Supervised by Dr. JET"
    ],
    "End confidently; thank the panel."
)

# Save file
prs.save("AI_Disease_Prediction_Presentation.pptx")
print("✅ Presentation generated successfully: AI_Disease_Prediction_Presentation.pptx")
